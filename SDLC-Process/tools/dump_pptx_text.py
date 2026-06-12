from pathlib import Path

from pptx import Presentation


def main() -> None:
    path = Path("SDLC-Process-v0.1.pptx")
    prs = Presentation(path)
    for i, slide in enumerate(prs.slides, start=1):
        print(f"\n=== Slide {i} ===")
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text = shape.text.strip()
                if text:
                    # Replace characters that cannot be encoded in the current console
                    safe = text.encode("cp1252", errors="replace").decode("cp1252")
                    print(safe)


if __name__ == "__main__":
    main()
