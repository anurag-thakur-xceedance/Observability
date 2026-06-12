from pathlib import Path

from pptx import Presentation
from pptx.util import Pt


def add_bullet_slide(prs: Presentation, title: str, bullets: list[str]) -> None:
    """Append a simple title+bullets slide."""
    # Try to find a layout with title and body; fall back to first layout.
    layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(layout)
    if slide.shapes.title is not None:
        slide.shapes.title.text = title

    body = None
    for shape in slide.placeholders:
        if hasattr(shape, "text_frame") and shape.placeholder_format.idx != 0:
            body = shape
            break

    if body is None:
        return

    tf = body.text_frame
    tf.clear()

    if not bullets:
        return

    p = tf.paragraphs[0]
    p.text = bullets[0]
    p.level = 0
    p.font.size = Pt(20)

    for bullet in bullets[1:]:
        p = tf.add_paragraph()
        p.text = bullet
        p.level = 0
        p.font.size = Pt(20)


def main() -> None:
    src = Path("SDLC-Process-v0.1.pptx")
    prs = Presentation(src)

    # Slide: SDLC In Your Repo & Pipeline
    add_bullet_slide(
        prs,
        "SDLC In Your Repo & Pipeline",
        [
            "Branching & PRs – Protected main/trunk branches with required PR reviews.",
            "Branching & PRs – Branch naming aligned to work items (e.g. feature/1234-short-title).",
            "Required Checks – Linting and unit tests must pass before merge.",
            "Required Checks – Secrets scan, SAST, and IaC checks wired as mandatory checks.",
            "Pipelines – Distinct stages: build → test → quality gates → package → deploy.",
            "Pipelines – Same pipeline definition across environments; configuration changes, not code.",
            "Artefacts – Versioned build artefacts and IaC templates stored centrally.",
            "Artefacts – Runbooks and change notes attached to releases.",
        ],
    )

    # Slide: What Good Looks Like (Metrics)
    add_bullet_slide(
        prs,
        "What Good Looks Like (Metrics)",
        [
            "Flow & Stability – Lead time for change trending down.",
            "Flow & Stability – Deployment frequency increasing where risk allows.",
            "Flow & Stability – Change fail rate and MTTR trending down.",
            "Quality – Defect escape rate to UAT/production decreasing.",
            "Quality – Coverage and linting warnings within agreed thresholds.",
            "Governance – Percentage of changes going through the full SDLC gates.",
            "Governance – Percentage of infrastructure changes done via IaC, not manual.",
            "Champion Role – Use these metrics in team reviews and retros; escalate systemic blockers.",
        ],
    )

    # Slide: Common Anti-Patterns To Watch For
    add_bullet_slide(
        prs,
        "Common Anti-Patterns To Watch For",
        [
            "Process Bypass – 'Quick fixes' merged directly to main without reviews or tests.",
            "Process Bypass – Manual environment changes never back-ported into IaC.",
            "Fake Quality Gates – Rubber-stamp reviews with no meaningful comments.",
            "Fake Quality Gates – Linting/tests configured but routinely ignored or overridden.",
            "Hero Dependency – One person who 'just knows' how to deploy or fix issues.",
            "Hero Dependency – Key knowledge not captured in runbooks, ADRs, or comments.",
            "Over-bureaucracy – Gates that add friction but no real risk reduction.",
            "Champion Response – Call these out early and propose automation/simplification.",
        ],
    )

    # Slide: Champion Enablement & Support
    add_bullet_slide(
        prs,
        "Champion Enablement & Support",
        [
            "Toolkit – SDLC Confluence space: phases, steps, and checklists.",
            "Toolkit – Templates for ADRs, test strategies, runbooks, and IaC/pipeline examples.",
            "Toolkit – Example repositories showing the 'golden path' implementations.",
            "Community – Developer Champions circles/guilds and regular syncs with CoE, Platform, Security.",
            "Support – Contact: <CoE Architecture DL / Teams channel> for guidance and escalation.",
            "Practice – Use structured feedback: what worked, what did not, and your proposal.",
            "Expectation – You are the feedback loop from delivery to architecture.",
            "Expectation – Help turn SDLC from a document into a living practice.",
        ],
    )

    # Try to overwrite the original; if it fails due to locking, write a v0.2 copy.
    try:
        prs.save(src)
    except PermissionError:
        backup = Path("SDLC-Process-v0.2.pptx")
        prs.save(backup)


if __name__ == "__main__":
    main()
