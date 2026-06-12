from pptx import Presentation
from pptx.util import Inches, Pt


def add_title_slide(prs: Presentation) -> None:
    slide_layout = prs.slide_layouts[0]  # Title slide
    slide = slide_layout and prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]

    title.text = "SDLC @ <Your Org>"
    subtitle.text = (
        "A Developer Champion's View\n"
        "Created by: CoE Architect / Enterprise Architect Team\n"
        "Audience: Developer Champions, Tech Leads, Principal Engineers"
    )


def add_bullet_slide(prs: Presentation, title_text: str, bullets: list[str]) -> None:
    slide_layout = prs.slide_layouts[1]  # Title and Content
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    body = slide.placeholders[1]

    title.text = title_text
    tf = body.text_frame
    tf.clear()

    if not bullets:
        return

    # First bullet
    p = tf.paragraphs[0]
    p.text = bullets[0]
    p.level = 0
    p.font.size = Pt(20)

    for bullet in bullets[1:]:
        p = tf.add_paragraph()
        p.text = bullet
        p.level = 0
        p.font.size = Pt(20)


def build_presentation() -> Presentation:
    prs = Presentation()

    # Slide 1: Title & Context
    add_title_slide(prs)

    # Slide 2: The "Why"
    add_bullet_slide(
        prs,
        "Why This SDLC Matters (To You)",
        [
            "Vision – One opinionated SDLC across the organisation.",
            "Vision – Built for modern cloud-native delivery (CI/CD, IaC, security-by-default).",
            "Vision – Lightweight enough for teams; strong enough for regulated clients.",
            "Problems – Fragmented ways of working across teams.",
            "Problems – Inconsistent quality gates (tests, reviews, security).",
            "Problems – Late surprises: infra drift, failed audits, broken releases.",
            "Problems – Hero culture: outcomes depend on individuals, not the system.",
            "Benefits – Clear, predictable path from idea to production.",
            "Benefits – Built-in hooks for testing, governance, and platform practices (IaC, secrets, observability).",
            "Benefits – Common language across squads, platform, security, and architecture.",
        ],
    )

    # Slide 3: Core SDLC Workflow
    add_bullet_slide(
        prs,
        "Core SDLC Workflow – Idea → Production",
        [
            "Discover & Shape – Capture demand, clarify business outcomes and scope; align on value, risks, and constraints.",
            "Design & Plan – Choose solution options and architecture; define NFRs, data, integration, security, and privacy.",
            "Develop – Agree branching and PR strategy; TDD/unit tests, linters, static analysis, comment-aware reasoning.",
            "Develop – Treat environments as code: IaC jobs and IaC testing; secrets scanning and secure-by-default patterns.",
            "Verify – Integration, contract, and E2E testing; performance, resilience, and security testing where relevant.",
            "Release & Operate – Promotion via pipelines; progressive delivery; observability, SLOs, and incident/change management.",
            "Learn & Improve – Post-implementation reviews, retros, and metrics feeding back into backlog and platform evolution.",
        ],
    )

    # Slide 4: Governance & Quality Gates
    add_bullet_slide(
        prs,
        "Key Governance & Quality Gates",
        [
            "Design / Plan – Architecture decisions recorded (ADRs or equivalent).",
            "Design / Plan – NFRs defined and testable (performance, security, compliance, availability).",
            "Design / Plan – Data and privacy impact understood with mitigations identified.",
            "Develop – Unit tests in place and passing; coverage thresholds agreed per team.",
            "Develop – Linting and static analysis clean enough to merge.",
            "Develop – Code review completed: functionality, readability, maintainability, and security impact.",
            "Develop – Meaningful comments, not drive-by approvals.",
            "Develop – IaC jobs and tests: infra changes via code, policy checks, drift detection, and cost impacts considered.",
            "Develop – Secrets scan clean: no hard-coded credentials or sensitive data in code or IaC.",
            "Verify / Release – Automated suites green for the scope being released.",
            "Verify / Release – Change is observable: logs, metrics, traces, alerts, and runbooks prepared.",
            "Verify / Release – Rollback or mitigation path defined and tested where feasible.",
            "Champions – Treat these as non-negotiable gates, not nice-to-haves.",
        ],
    )

    # Slide 5: Org-Wide Implementation Plan
    add_bullet_slide(
        prs,
        "Org-Wide Implementation Plan",
        [
            "Standardise – Confluence as the single source of truth for SDLC steps and checklists.",
            "Standardise – Phase/step-level guidance (e.g. 'Develop – Code Review', 'Develop – Secrets Scan').",
            "Standardise – Templates for ADRs, test strategies, runbooks, and IaC/pipeline patterns.",
            "Tooling – CI/CD pipelines aligned to SDLC phases: linting, tests, SAST, IaC checks, secrets, packaging, deployment.",
            "Tooling – Default repo scaffolds with branching model, quality gate configs, and example pipelines/policies.",
            "Rollout – Wave 1: pilot with selected product teams and platform.",
            "Rollout – Wave 2: expand to all new projects and major change initiatives.",
            "Rollout – Wave 3: retrofit critical legacy services where risk is highest.",
            "Feedback – Regular syncs between CoE Architects, Developer Champions, Platform, and Security.",
            "Feedback – Data-driven adjustments using lead time, deployment frequency, MTTR, change fail rate, and defect escape rate.",
        ],
    )

    # Slide 6: Champion's Call to Action
    add_bullet_slide(
        prs,
        "Developer Champion – Call to Action",
        [
            "Adoption – Map your team's current way of working to the SDLC steps and identify gaps.",
            "Adoption – Ensure core gates (tests, review, IaC, secrets, observability) are enforced.",
            "Raising the bar – Challenge low-value, 'tick-box' gates; automate manual, repeatable checks.",
            "Raising the bar – Make the right way the easy way for your engineers.",
            "Feedback – Capture friction points, duplications, and missing guidance.",
            "Feedback – Bring real examples back to the CoE Architect team.",
            "Next steps – Pick one active initiative and align its workflow to the SDLC phases.",
            "Next steps – Join or start a Developer Champions circle for your domain or tribe.",
            "Next steps – Use the SDLC Confluence space as your working reference and contribute improvements.",
        ],
    )

    return prs


def main() -> None:
    prs = build_presentation()
    # Use a v2 filename to avoid issues if the original file is open
    output_path = "SDLC-Developer-Champions-v2.pptx"
    prs.save(output_path)


if __name__ == "__main__":
    main()
